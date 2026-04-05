import os
import sqlite3
import json
import pandas as pd
from pptx import Presentation
from groq import Groq
from dotenv import dotenv_values

env_vars = dotenv_values(".env")
GroqAPIKey = env_vars.get("GroqAPIKey")

def GetStructuredData(prompt, format_type):
    client = Groq(api_key=GroqAPIKey)
    
    if format_type == "ppt":
        sys_msg = "Output a valid JSON array of objects. Each object represents a slide with keys: 'title' and 'content'."
    elif format_type == "excel":
        sys_msg = "Output valid JSON representing tabular data. It must be a list of dictionaries."
    else:
        sys_msg = "Output valid JSON."

    try:
        # Simplified string to prevent termination errors
        full_user_query = f"Create data for: {prompt}. OUTPUT STRICTLY JSON ONLY."
        
        completion = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {"role": "system", "content": sys_msg},
                {"role": "user", "content": full_user_query}
            ],
            temperature=0.2
        )
        
        raw = completion.choices[0].message.content.strip()
        
        # Ultra-safe cleanup
        raw = raw.replace("```json", "").replace("```", "").strip()
        
        return json.loads(raw)
    except Exception as e:
        print(f"!! Data Structure Error: {e}")
        return None

def CreatePresentation(topic, filename="Presentation.pptx"):
    save_dir = os.path.join(os.getcwd(), "Workspace", "Presentations")
    os.makedirs(save_dir, exist_ok=True)
    file_path = os.path.join(save_dir, filename)

    print(f">> [WORKSPACE]: Drafting presentation on {topic}...")
    slides_data = GetStructuredData(topic, "ppt")
    
    if not slides_data: 
        return "I failed to structure the presentation data, sir."

    prs = Presentation()
    for slide_info in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = str(slide_info.get("title", "Slide"))
        slide.placeholders[1].text = str(slide_info.get("content", ""))

    prs.save(file_path)
    return f"Presentation created successfully. Saved to {file_path}"

def CreateExcel(topic, filename="DataSheet.xlsx"):
    save_dir = os.path.join(os.getcwd(), "Workspace", "Spreadsheets")
    os.makedirs(save_dir, exist_ok=True)
    file_path = os.path.join(save_dir, filename)

    print(f">> [WORKSPACE]: Compiling dataset for {topic}...")
    excel_data = GetStructuredData(topic, "excel")
    
    if not excel_data: 
        return "I failed to compile the spreadsheet data, sir."

    df = pd.DataFrame(excel_data)
    df.to_excel(file_path, index=False)
    return f"Spreadsheet compiled successfully. Saved to {file_path}"

def CreateDatabase(topic, db_name="JarvisData.db"):
    save_dir = os.path.join(os.getcwd(), "Workspace", "Databases")
    os.makedirs(save_dir, exist_ok=True)
    file_path = os.path.join(save_dir, db_name)

    print(f">> [WORKSPACE]: Architecting Database for {topic}...")
    db_data = GetStructuredData(topic, "excel") 
    
    if not db_data: 
        return "I failed to architect the database schema, sir."

    df = pd.DataFrame(db_data)
    conn = sqlite3.connect(file_path)
    table_name = "".join(e for e in topic if e.isalnum()) or "DataTable"
    df.to_sql(table_name, conn, if_exists='replace', index=False)
    conn.close()
    
    return f"Database architected successfully. Saved to {file_path}"